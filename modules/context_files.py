# modules/context_files.py
# ─────────────────────────────────────────────────────────────────────────────
# Extração de texto de arquivos de referência do contexto.
#
# Formatos suportados:
#   .html / .htm  — strip de tags via html.parser (stdlib, sem dependência extra)
#   .pptx         — extração de slides + notas via python-pptx
#   .pdf          — extração via pypdf (já em requirements.txt)
#   .txt / .md    — decodificação direta UTF-8
#
# Segurança HTML: apenas o texto é extraído — scripts, iframes e atributos
# são descartados pelo parser. Nenhum HTML é renderizado no app.
# ─────────────────────────────────────────────────────────────────────────────

from __future__ import annotations

import io
import re

SUPPORTED_EXTENSIONS = {".html", ".htm", ".pptx", ".pdf", ".txt", ".md"}
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB


def extract_text(filename: str, file_bytes: bytes) -> str:
    """
    Extrai texto legível de um arquivo de referência.

    Args:
        filename:   nome do arquivo (usado para detectar o formato pela extensão)
        file_bytes: conteúdo binário do arquivo

    Returns:
        Texto extraído como string UTF-8, pronto para injeção no prompt.
        Retorna string vazia se a extração falhar.
    """
    ext = _ext(filename)
    try:
        if ext in (".html", ".htm"):
            return _extract_html(file_bytes)
        if ext == ".pptx":
            return _extract_pptx(file_bytes)
        if ext == ".pdf":
            return _extract_pdf(file_bytes)
        if ext in (".txt", ".md"):
            return file_bytes.decode("utf-8", errors="replace").strip()
    except Exception:
        return ""
    return ""


def _ext(filename: str) -> str:
    from pathlib import Path
    return Path(filename).suffix.lower()


# ── HTML ──────────────────────────────────────────────────────────────────────

def _extract_html(raw: bytes) -> str:
    """Extrai texto de HTML usando html.parser da stdlib — zero dependências."""
    from html.parser import HTMLParser

    class _TextExtractor(HTMLParser):
        _SKIP = {"script", "style", "head", "meta", "link", "noscript"}

        def __init__(self):
            super().__init__()
            self._parts: list[str] = []
            self._skip_depth = 0

        def handle_starttag(self, tag, attrs):
            if tag.lower() in self._SKIP:
                self._skip_depth += 1

        def handle_endtag(self, tag):
            if tag.lower() in self._SKIP and self._skip_depth > 0:
                self._skip_depth -= 1

        def handle_data(self, data):
            if self._skip_depth == 0:
                stripped = data.strip()
                if stripped:
                    self._parts.append(stripped)

        def get_text(self) -> str:
            return "\n".join(self._parts)

    try:
        html_str = raw.decode("utf-8", errors="replace")
    except Exception:
        return ""

    parser = _TextExtractor()
    parser.feed(html_str)
    text = parser.get_text()
    # Collapse excessive blank lines
    return re.sub(r"\n{3,}", "\n\n", text).strip()


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _extract_pptx(raw: bytes) -> str:
    """Extrai texto de apresentação PowerPoint (.pptx) via python-pptx."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return "(python-pptx não instalado — não foi possível extrair texto do PPTX)"

    prs = Presentation(io.BytesIO(raw))
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = [f"### Slide {i}"]

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(r.text for r in para.runs).strip()
                if line:
                    slide_parts.append(line)

        # Speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_frame.text.strip() if notes_frame else ""
            if notes_text:
                slide_parts.append(f"*Notas:* {notes_text}")

        if len(slide_parts) > 1:
            parts.append("\n".join(slide_parts))

    return "\n\n".join(parts).strip()


# ── PDF ───────────────────────────────────────────────────────────────────────

def _extract_pdf(raw: bytes) -> str:
    """Extrai texto de PDF via pypdf."""
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError:
        return "(pypdf não instalado — não foi possível extrair texto do PDF)"

    reader = PdfReader(io.BytesIO(raw))
    pages: list[str] = []
    for i, page in enumerate(reader.pages, 1):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append(f"### Página {i}\n{text}")

    return "\n\n".join(pages).strip()
